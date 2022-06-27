# 使用python操作ppt需要使用python-pptx 第三方库
# 安装：pip install python-pptx
import pptx
from pptx import Presentation

# 创建幻灯片
pres = Presentation()

# 选择母版添加一页
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)

# 获取标题栏和副标题栏
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 编辑标题和副标题
title.text = "welcome to python"
subtitle.text = "Life is sort ,I use Python"

# 保存幻灯片
pres.save("./data/test.pptx")
